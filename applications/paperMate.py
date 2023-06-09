# This application provides fully or partially automated assistance for writing papers and using a system.
# The company's confidential information has been erased or modified in the codes
# By using JSON files, important information can be kept hidden from view and be changed easily without modifying the codes.
# Users are also able to modify some information in the application's configuration settings(설정)

import tkinter as tk
from tkinter import ttk
from tkinter import messagebox
from tkinter import filedialog
import locale
import os
import json
import datetime


# Set locale for timestamp and formatting price
locale.setlocale(locale.LC_ALL, '')

def timeStamp():
    now = datetime.datetime.now()

    return now.strftime("%Y-%m-%d %H:%M:%S")

def loadData():
    with open('./pmdata/data.json', 'r') as file:
        global data
        data = json.load(file)

loadData()
    
def saveData(data):
    with open('./pmdata/data.json', 'w') as file:
        # write the dictionary to the file
        json.dump(data, file)
        
def saveDb(record):
    with open('db.json', 'r+') as file:
        db = json.load(file)
        db.append(record)
        file.seek(0)
        json.dump(db, file)
        
    return len(db)
    
#Function for changing numerical prices into their Korean chracter equivalents
def priceToHangul(num):
    numPrice = int(num)
    unitLi = []
    unitHan = ['천백십 ', '만', '억']
    numHan = '일이삼사오육칠팔구'
    priceHan = []
    while numPrice >= 1:
        unitLi.append('{:0>4}'.format(str(numPrice%10000)))
        numPrice //= 10000
    for i in range(len(unitLi)):
        temPrice = "".join([numHan[int(a)-1]+b for a, b in zip(unitLi[i], unitHan[0]) if int(a)])
        if i:
            temPrice += unitHan[i]
        priceHan.append(temPrice)
    priceHan.reverse()
    result = "".join(priceHan).replace(" ", "")
    
    return result

# Define a class for entry with placeholder
class PlaceholderEntry(tk.Entry):
    def __init__(self, master=None, placeholder="Enter text...", color='grey'):
        super().__init__(master)
        self.placeholder = placeholder
        self.placeholder_color = color
        self.default_fg_color = self['fg']
        self.bind("<FocusIn>", self._clear_placeholder)
        self.bind("<FocusOut>", self._add_placeholder)
        self._add_placeholder()

    def _clear_placeholder(self, event):
        if self['fg'] == self.placeholder_color:
            self.delete(0, "end")
            self.config(fg=self.default_fg_color)

    def _add_placeholder(self, event=None):
        if not self.get():
            self.insert(0, self.placeholder)
            self.config(fg=self.placeholder_color)

# GUI for the app            
class windows_tkinter:
    def __init__(self, window):
        self.window = window
        self.window.title("paperMate")
        self.window.iconbitmap('./icon/icon.ico')
        self.icon_image = tk.PhotoImage(file='./icon/icon.png')
        self.window.iconphoto(True, self.icon_image)
        self.window.geometry(f"320x400+{screen_width//2 - 160}+{screen_height//2 - 270}")
        self.window.resizable(0, 0)
        self.window.focus_force()

        self.__main__()
        
    def close(self):
        self.window.quit()
        self.window.destroy()
       
    
    def startWriting(self, info, inputId, inputPw):
        try:
            #Import selenium library for automation
            from selenium import webdriver
            from selenium.webdriver.chrome.options import Options
            from selenium.webdriver.common.by import By
            from selenium.webdriver.common.keys import Keys
            from selenium.webdriver.support.select import Select
            from selenium.webdriver.support.ui import WebDriverWait
            from selenium.webdriver.support import expected_conditions as EC
            
            #Define variables for writing contract
            title = info['title']
            tests = info['tests']
            price = info['price']
            date_contract = info['date_contract']
            date_end = info['date_end']
            company = info['company']
            loginId = inputId
            loginPw = inputPw

            price = int(price)
            payment = int(price/2)

            #initializing browser
            chrome_options = Options()
            driver = webdriver.Chrome(options=chrome_options)
            #get to the portal
            url = 'https://' 
            driver.get(url)

            driver.implicitly_wait(0.5)

            #user login
            userId = driver.find_element(by=By.NAME, value="user_id")
            userPw = driver.find_element(by=By.NAME, value="pswd")
            userId.send_keys(loginId)
            userPw.send_keys(loginPw)

            loginBtn = driver.find_element(by=By.CLASS_NAME, value="btn-login")
            loginBtn.click()
            driver.implicitly_wait(1)

            #move to 법무시스템
            glas = "https://"
            driver.get(glas)
            driver.implicitly_wait(1)

            #open 표준계약서 검토
            element = driver.find_element(by=By.XPATH, value="/html/body/div[3]/ul[1]/li[3]")
            driver.execute_script("arguments[0].click();", element)
            driver.implicitly_wait(1)

            #To find exact window and switch to it
            numWindows = len(driver.window_handles)
            for i in range(1, numWindows):
                driver.switch_to.window(driver.window_handles[i])
                try:
                    if driver.find_element(By.CSS_SELECTOR, '.content_name').text == "표준계약 의뢰서":
                        break
                except:
                        pass

            #Input basic information
            driver.find_element(by=By.CSS_SELECTOR, value="#lms_pati_gyeyag_tit").send_keys(f"{title} 계약")
            Select(driver.find_element(by=By.CSS_SELECTOR, value="#lms_pati_hoesa_cod")).select_by_value("L")
            driver.find_element(by=By.CSS_SELECTOR, value="#rdo_ins_type0").click()
            driver.find_element(by=By.CSS_SELECTOR, value="#lms_pati_gyeyag_geumaeg").send_keys(price)
            driver.find_element(by=By.CSS_SELECTOR, value="#lms_pati_teugsu_jogeon_yn0").click()
            driver.find_element(by=By.CSS_SELECTOR, value="#rdo_gaein_weetak_yn0").click()
            
            #Push '계약 상대방 검색' button
            searchCounterpart = driver.find_element(by=By.CSS_SELECTOR, value="#a_sch_ctp")
            driver.execute_script("arguments[0].click();", searchCounterpart)
            driver.implicitly_wait(1)
            
            #Search the company
            companyCode = data['companies'][company]['code']
            driver.switch_to.window(driver.window_handles[-1])
            searchKeyword = driver.find_element(by=By.XPATH, value="/html/body/div[2]/form/table[1]/tbody/tr[2]/td[2]/table/tbody/tr[1]/td[1]/input")
            searchKeyword.send_keys(company)
            searchKeyword.send_keys(Keys.ENTER)
            wait = WebDriverWait(driver, 10)
            wait1 = wait.until(EC.presence_of_element_located((By.CSS_SELECTOR, f'#{companyCode} .radio')))
            driver.find_element(by=By.CSS_SELECTOR, value=f"#{companyCode} .radio").click()
            wait2 = wait.until(EC.text_to_be_present_in_element((By.CSS_SELECTOR, '#insert_td'), '거래가능'))
            returnCounterpart = driver.find_element(by=By.XPATH, value="/html/body/div[2]/form/div[1]/input[2]")
            driver.execute_script("arguments[0].click();", returnCounterpart)
            driver.implicitly_wait(3)
            
            #Back to the window for basic information(use variable 'i' from above)
            driver.switch_to.window(driver.window_handles[i])
            driver.find_element(By.CSS_SELECTOR, '#span_sch_std_ct').click()
            driver.implicitly_wait(3)

            #Get to new window for contract
            driver.switch_to.window(driver.window_handles[-1])
            searchContract = driver.find_element(By.XPATH, '/html/body/div[2]/div/form/table[1]/tbody/tr[2]/td[2]/table/tbody/tr/td[2]/input')
            searchContract.send_keys('업무 계약')
            searchContract.send_keys(Keys.ENTER)
            driver.implicitly_wait(3)
            driver.find_element(By.XPATH, "/html/body/div[2]/div/form/table[4]/tbody/tr[1]/td[1]/input[1]").click()
            driver.implicitly_wait(3)

            driver.switch_to.window(driver.window_handles[-1])
            
            #contract_constant(can be changed in '설정'), last modified 2023-03-27 
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val0').send_keys(data['common']['nameLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val1').send_keys(data['companies'][company]['nameCompany'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val2').send_keys(data['companies'][company]['nameResearchLead'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val6').send_keys(data['common']['dayPayment'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val8').send_keys(data['common']['dayPayment'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val11').send_keys(data['common']['adressLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val12').send_keys(data['common']['ceoLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val13').send_keys(data['companies'][company]['nameCompany2'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val14').send_keys(data['companies'][company]['adressCompany'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val15').send_keys(data['companies'][company]['ceoCompany'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val16').send_keys(data['companies'][company]['adressResearchLead'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val17').send_keys(data['companies'][company]['nameResearchLead'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val18').send_keys(data['common']['nameLG2'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val19').send_keys(data['common']['adressLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val20').send_keys(data['common']['businessLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val21').send_keys(data['common']['regNumLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val22').send_keys(data['common']['personLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val23').send_keys(data['common']['personTeamLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val24').send_keys(data['common']['personAddressLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val25').send_keys(data['common']['personContactLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val26').send_keys(data['common']['personEmailLG'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val28').send_keys(data['common']['howToUse'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val29').send_keys(data['common']['howToStore'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val30').send_keys(data['common']['howToUse'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val31').send_keys(data['common']['howToUse'])
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val33').send_keys(data['common']['nameLG2'])

            #contract_variables 
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val3').send_keys(title)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val4').send_keys(price)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val5').send_keys(payment)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val7').send_keys(payment)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val10').send_keys(date_contract)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val9').send_keys(date_end)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val32').send_keys(date_contract)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val34').send_keys(date_contract)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val35').send_keys(date_contract)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val36').send_keys(date_contract)
            driver.find_element(By.CSS_SELECTOR, '#std_ct_param_val27').send_keys(tests)

            #Show preview of contract 
            driver.find_element(By.CSS_SELECTOR, '#btn_apply').click()
            
            info['timestamp'] = timeStamp()
            info['origin'] = 'contract'
            saveDb(info)
            messagebox.showinfo('완료', '입력 내용 검토 후 계약 검토를 완료하세요.\n완료 후에는 자동화 브라우저 종료해주세요', parent = self.tkWindowContract)

            #To keep browser open
            while(driver.window_handles):
                pass

            driver.quit()
            
        except Exception as e:
            driver.quit()
            if '아이디 또는 비밀번호가 다릅니다' in str(e):
                messagebox.showwarning('Login failed', f'''비밀번호가 잘못 되었어요. 다시 확인해주세요
                
                {str(e)}''', parent=self.tkWindowContract)
            else:                                      
                messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.tkWindowContract)
            
        
    def writeContract(self, info):
        try:
            for child in self.tkWindowContract.winfo_children():
                if isinstance(child, tk.Toplevel):
                    if child.title() == '로그인 정보 입력':
                        child.destroy()
                    
            self.info = info

            self.tkWindowLogin = tk.Toplevel(self.tkWindowContract)
            self.tkWindowLogin.geometry(f"300x220+{screen_width//2 - 150}+{screen_height//2 - 250}")
            self.tkWindowLogin.title("로그인 정보 입력")
            self.tkWindowLogin.resizable(0, 0)
            self.tkWindowLogin.focus_force()

            loginLb = tk.Label(self.tkWindowLogin, text = "")
            loginLb.grid(row=1, column=1, pady = 10)
            lbId = tk.Label(self.tkWindowLogin, text = "ID")
            lbId.grid(row=2, column=2, pady = 10, padx = (9, 10))
            setId = tk.StringVar()
            #Set saved ID in entry if saveId == 1
            if data['myInfo']['saveId']:
                setId.set(data['myInfo']['userId'])
            entryId = tk.Entry(self.tkWindowLogin, textvariable=setId)
            entryId.grid(row=2, column = 3)
            lbPw = tk.Label(self.tkWindowLogin, text = "Password")
            lbPw.grid(row=3, column=2, pady = (2, 0), padx = (9, 10))
            entryPw = tk.Entry(self.tkWindowLogin, show = '*')
            entryPw.grid(row=3, column = 3, pady = (2, 0))
            tk.Label(self.tkWindowLogin, text = "").grid(row=4, column=1, pady=(0, 5))
               
            rememberId = tk.BooleanVar()
            rememberIdBtn = tk.Checkbutton(self.tkWindowLogin, text="ID 기억하기", variable=rememberId)
            rememberIdBtn.grid(row = 5, column = 3)
            if data['myInfo']['saveId']:
                rememberIdBtn.select()

            def validateLogin(info, entryId, entryPw):
                import re
                #To validate if there's any korean character in inputted password
                korean_regex = re.compile("[\u3131-\u3163\uac00-\ud7a3]+")
                bg_color = self.tkWindowLogin.cget('background')
                self.Id = entryId.get()
                self.Pw = entryPw.get()
                valiCnt = 0
                idWarning=tk.Message(self.tkWindowLogin, text="!", 
                                            width=10, padx = 10, foreground='red')
                pwWarning=tk.Message(self.tkWindowLogin, text="!", 
                                            width=10, padx = 10, foreground='red')
                pwWarningHan = tk.Message(self.tkWindowLogin, text = "비밀번호에 한글 포함", 
                                          width=150, foreground = 'red')
                
                #Validate if id and password were inputted correctly
                if not self.Id or korean_regex.search(self.Id):
                    idWarning.grid(row = 2, column = 4, padx = (5, 0))
                else:
                    valiCnt += 1
                    tk.Message(self.tkWindowLogin, text=" ", 
                                            width=100, padx = 10).grid(row=2, column=4)
                if not self.Pw:
                    pwWarning.grid(row = 3, column = 4, pady = (0, 0), padx = (5, 0))
                elif korean_regex.search(self.Pw):
                    pwWarning.grid(row = 3, column = 4, pady = (0, 0), padx = (5, 0))
                    pwWarningHan.grid(row=4, column= 3, columnspan=2, sticky = 'w')
                    #Warning message will flicker
                    f = 1
                    while f < 8:
                        if f%2 == 1:
                            self.tkWindowLogin.after(180*f, lambda: pwWarningHan.config(foreground=bg_color))
                        else:
                            self.tkWindowLogin.after(180*f, lambda: pwWarningHan.config(foreground='red'))
                        f += 1
                else:
                    valiCnt += 1
                    tk.Message(self.tkWindowLogin, text=" ", 
                                            width=100, padx = 10).grid(row=3, column=4)

                if valiCnt == 2:
                    #If user checked 'ID 기억하기', change status fo 'saveId' and save inputted ID to data.json
                    if rememberId.get():
                        if data['myInfo']['saveId'] != 1 or data['myInfo']['userId'] != self.Id:
                            data['myInfo']['saveId'] = 1
                            data['myInfo']['userId'] = self.Id
                            saveData(data)
                            loadData()
                    else:
                        if data['myInfo']['saveId'] != 0:
                            data['myInfo']['saveId'] = 0
                            saveData(data)
                            loadData()
                    
                    self.tkWindowLogin.destroy()
                    self.startWriting(info, self.Id, self.Pw)
                    


            entryPw.bind('<Return>', lambda event: validateLogin(self.info, entryId, entryPw))

            loginBtn = tk.Button(self.tkWindowLogin, overrelief="solid", 
                           command = lambda: validateLogin(self.info, entryId, entryPw), 
                           text = "확인")
            loginBtn.grid(row=6, column=3, ipadx = 20, pady = 15)    
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.tkWindowContract)
   
    def writeChu(self, info):
        try:
            from pptx import Presentation, util

            #Define variables for writing 추진사유서
            tests = info['tests']
            testsList = tests.split(',')
            if len(testsList) == 1:
                tests2 = tests
            else:
                tests2 = testsList[0]+f' 외 {len(testsList)-1}건'
            pdct = info['pdct']
            company = info['company']
            price = info['price']
            priceHangul = priceToHangul(price)
            
            # Load a preset PowerPoint file
            prs = Presentation(f'./pmdata/{company}.pptx')

            chuTitle = prs.slides[0].shapes[1]
            chuTitle.text = chuTitle.text.replace("{pdct}", pdct)
            chuTitle.text = chuTitle.text.replace("{tests2}", tests2)
            chuTitle.text_frame.paragraphs[0].font.size = util.Pt(10)

            chuTb = prs.slides[0].shapes[2].table
            chuTbCell1 = chuTb.cell(0,1)
            chuTbCell1.text = chuTbCell1.text.replace('{pdct}', pdct)
            chuTbCell1.text = chuTbCell1.text.replace('{tests}', tests)
            chuTbCell1.text_frame.paragraphs[0].font.size = util.Pt(11)

            chuTbCell2 = chuTb.cell(3,1)
            chuTbCell2.text = chuTbCell2.text.replace('{pdct}', pdct)
            chuTbCell2.text = chuTbCell2.text.replace('{tests}', tests)
            chuTbCell2.text_frame.paragraphs[0].font.size = util.Pt(11)

            chuTbCell3 = chuTb.cell(4,1)
            chuTbCell3.text = chuTbCell3.text.replace('{priceHangul}', priceHangul)
            chuTbCell3.text_frame.paragraphs[0].font.size = util.Pt(11)
            
            #Ask users where to save the file
            dir_path = filedialog.askdirectory(parent = self.tkWindowContract)
            if dir_path:
                prs.save(dir_path + f'/추진사유서_{tests2}_{pdct}({company}).pptx')
                
                info['timestamp'] = timeStamp()
                info['origin'] = 'chu'
                saveDb(info)
                messagebox.showinfo("안내", '추진사유서 저장되었습니다.', parent=self.tkWindowContract)
            else:
                pass
                
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.tkWindowContract)

            
    def writeMail(self, info, userName):
        try:
            import win32com.client as win32

            #Define variables for writing mail
            price = int(info['price'])
            halfPrice = locale.currency(price//2, grouping=True)
            price = locale.currency(price, grouping=True)

            tests = info['tests']
            testsList = tests.split(',')
            if len(testsList) == 1:
                tests2 = tests
            else:
                tests2 = testsList[0]+f' 외 {len(testsList)-1}건'
            pdct = info['pdct']
            company = info['company']
            date_contract = info['date_contract']
            date_end = info['date_end']

            # Get the current selection in the active Word document
            word = win32.Dispatch("Word.Application")
            word.Visible = False
            path_mailPum = 'pmdata/mailPum.docx'
            abspath_mailPum = os.path.abspath(path_mailPum)
            doc = word.Documents.Open(abspath_mailPum)
            story_range = doc.StoryRanges(win32.constants.wdMainTextStory)

            change1 = story_range.Find.Execute('{userName}', False, False, False, False, False, True, 1, True, userName, 2)
            change2 = story_range.Find.Execute('{tests2}', False, False, False, False, False, True, 1, True, tests2, 2)
            change3 = story_range.Find.Execute('{pdct}', False, False, False, False, False, True, 1, True, pdct, 2)
            change4 = story_range.Find.Execute('{price}', False, False, False, False, False, True, 1, True, price, 2)
            change5 = story_range.Find.Execute('{tests}', False, False, False, False, False, True, 1, True, tests, 2)
            change6 = story_range.Find.Execute('{company}', False, False, False, False, False, True, 1, True, company, 2)
            change7 = story_range.Find.Execute('{date_contract}', False, False, False, False, False, True, 1, True, date_contract, 2)
            change8 = story_range.Find.Execute('{date_end}', False, False, False, False, False, True, 1, True, date_end, 2)
            change9 = story_range.Find.Execute('{halfPrice}', False, False, False, False, False, True, 1, True, halfPrice, 2)

            doc.Content.WholeStory()
            doc.Content.Copy()

            #To keep the original word file unchanged 
            doc.Saved = True
            
            doc.Close()
            
            info['timestamp'] = timeStamp()
            info['origin'] = 'mail'
            saveDb(info)
            messagebox.showinfo("안내", '메일 내용이 복사되었습니다', parent=self.tkWindowContract)

            #If there's no word file opened, quit word application
            otherWord = win32.GetActiveObject("Word.Application")
            if not otherWord.Documents:
                word.Quit()
            else:
                pass
                
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.tkWindowContract)
        
    def askName(self, info):
        try:
            #Remove a window if there exist one already 
            for child in self.tkWindowContract.winfo_children():
                if isinstance(child, tk.Toplevel):
                    if child.title() == '이름 작성':
                        child.destroy()
                    
            self.info = info

            self.tkWindowAskName = tk.Toplevel(self.tkWindowContract)
            self.tkWindowAskName.geometry(f"250x150+{screen_width//2 - 125}+{screen_height//2 - 90}")
            self.tkWindowAskName.title("이름 작성")
            self.tkWindowAskName.resizable(0, 0)
            self.tkWindowAskName.focus_force()


            askNameLb = tk.Label(self.tkWindowAskName, text = "")
            askNameLb.grid(row=1, column=1, pady = 5)
            lbName = tk.Label(self.tkWindowAskName, text = "이름")
            lbName.grid(row=2, column=2, pady = 10, padx = 15)
            setName = tk.StringVar()
            
            #Set saved name in name entry if saveName == 1
            if data['myInfo']['saveName']:
                setName.set(data['myInfo']['userName'])
            entryName = tk.Entry(self.tkWindowAskName, textvariable=setName)
            entryName.grid(row=2, column = 3)

            rememberName = tk.BooleanVar()
            rememberNameBtn = tk.Checkbutton(self.tkWindowAskName, text="이름 기억하기", variable=rememberName)
            rememberNameBtn.grid(row = 4, column = 3)
            if data['myInfo']['saveName']:
                rememberNameBtn.select()

            def validateName(info, entryName):
                userName = entryName.get()
                nameWarning=tk.Message(self.tkWindowAskName, text="!", 
                                            width=10, padx = 10, foreground='red')

                if not userName:
                    nameWarning.grid(row = 2, column = 4)
                else:
                    tk.Message(self.tkWindowAskName, text=" ", 
                                            width=100, padx = 10).grid(row=2, column=4)

                    if rememberName.get():
                        if data['myInfo']['saveName'] != 1 or data['myInfo']['userName'] != userName:
                            data['myInfo']['saveName'] = 1
                            data['myInfo']['userName'] = userName
                            saveData(data)
                            loadData()
                    else:
                        if data['myInfo']['saveName'] != 0:
                            data['myInfo']['saveName'] = 0
                            saveData(data)
                            loadData()
    
                    self.tkWindowAskName.destroy()
                    self.writeMail(info, userName)
                    
            entryName.bind('<Return>', lambda event: validateName(self.info, entryName))
            mailBtn = tk.Button(self.tkWindowAskName, overrelief="solid", 
                           command = lambda: validateName(self.info, entryName), 
                           text = "확인")
            mailBtn.grid(row=5, column=3, ipadx = 20, pady = 15)    
            
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.tkWindowContract)

    
    def startContract(self):
        try:
            for child in self.window.winfo_children():
                if isinstance(child, tk.Toplevel):
                    if child.title() == '새 업무 시작':
                        child.destroy()

            self.tkWindowContract = tk.Toplevel(self.window)
            self.tkWindowContract.geometry(f"520x450+{screen_width//2 - 250}+{screen_height//2 - 250}")
            self.tkWindowContract.title("새 업무 시작")
            self.tkWindowContract.resizable(0, 0)
            self.tkWindowContract.focus_force()
            
            conLb = tk.Label(self.tkWindowContract, text = "")
            conLb.grid(row=1, column=1)
            conLb1 = tk.Label(self.tkWindowContract, text = "")
            conLb1.grid(row=2, column=1,  padx = (20, 30), pady = 10)
            conLbTitle = tk.Label(self.tkWindowContract, text = "계약명")
            conLbTitle.grid(row=2, column=2, padx = 15)
            conEnTitle = PlaceholderEntry(self.tkWindowContract, placeholder = "OO 계약")
            conEnTitle.grid(row=2, column = 3, ipadx = 40)
            conLbPdct = tk.Label(self.tkWindowContract, text = "제품명")
            conLbPdct.grid(row=3, column=2, pady = 10)
            conEnPdct = PlaceholderEntry(self.tkWindowContract, placeholder = "")
            conEnPdct.grid(row=3, column = 3, ipadx = 40)
            conLb2 = tk.Label(self.tkWindowContract, text = "")
#             conLb2.grid(row=2, column=1,  padx = 30)
            conLbCom = tk.Label(self.tkWindowContract, text = "계약 항목")
            conLbCom.grid(row=4, column=2, pady = 10)
            conEnCom = PlaceholderEntry(self.tkWindowContract, placeholder = "")
            conEnCom.grid(row=4, column = 3, ipadx = 40)
            conLb3 = tk.Label(self.tkWindowContract, text = "")
#             conLb3.grid(row=3, column=1,  padx = 30)
            conLbPrice = tk.Label(self.tkWindowContract, text = "계약 금액")
            conLbPrice.grid(row=5, column=2, pady = 10)
            conEnPrice = PlaceholderEntry(self.tkWindowContract, placeholder = "숫자만 입력")
            conEnPrice.grid(row=5, column = 3, ipadx = 40)
            conLb4 = tk.Label(self.tkWindowContract, text = "")
#             conLb4.grid(row=4, column=1,  padx = 30)
            conLbStartDate = tk.Label(self.tkWindowContract, text = "계약일")
            conLbStartDate.grid(row=6, column=2, pady = 10)
            conEnStartDate = PlaceholderEntry(self.tkWindowContract, placeholder = "YYYYMMDD")
            conEnStartDate.grid(row=6, column = 3, ipadx = 40)
            conLb5 = tk.Label(self.tkWindowContract, text = "")
#             conLb5.grid(row=5, column=1,  padx = 30)
            conLbEndDate = tk.Label(self.tkWindowContract, text = "종료일")
            conLbEndDate.grid(row=7, column=2, pady = 10)
            conEnEndDate = PlaceholderEntry(self.tkWindowContract, placeholder = "YYYYMMDD")
            conEnEndDate.grid(row=7, column = 3, ipadx = 40)

            #Get list of companies from data.json then make a combobox
            companies = list(data['companies'].keys())
            selectedCompany = tk.StringVar()
            comboboxCompanies = ttk.Combobox(self.tkWindowContract, 
                                             value = companies,
                                            textvariable = selectedCompany, state = 'readonly')
            conLbCompanies = tk.Label(self.tkWindowContract, text = "계약 상대")
            conLbCompanies.grid(row=8, column=2, pady = 10)
            comboboxCompanies.grid(row = 8, column = 3, ipadx = 30)

            #Validate inputs
            def goContract(title, pdct, tests, price, date_1, date_2, company, btn):
                self.title = title.get()
                self.pdct = pdct.get()
                self.tests = tests.get()
                self.price = price.get()
                self.date_1 = date_1.get()
                self.date_2 = date_2.get()
                self.company = company.get()
                message1=tk.Message(self.tkWindowContract, text="여기 잘못 쓰셨네요", 
                                            width=200, padx = 10, foreground='red')
                message2=tk.Message(self.tkWindowContract, text="여기 잘못 쓰셨네요", 
                                            width=200, padx = 10, foreground='red')
                message3=tk.Message(self.tkWindowContract, text="여기 잘못 쓰셨네요", 
                                            width=200, padx = 10, foreground='red')
                message4=tk.Message(self.tkWindowContract, text="여기 잘못 쓰셨네요", 
                                            width=200, padx = 10, foreground='red')
                message5=tk.Message(self.tkWindowContract, text="여기 잘못 쓰셨네요", 
                                            width=200, padx = 10, foreground='red')
                message6=tk.Message(self.tkWindowContract, text="여기 잘못 쓰셨네요", 
                                            width=200, padx = 10, foreground='red')
                message7=tk.Message(self.tkWindowContract, text="여기 안 고르셨네요", 
                                            width=200, padx = 10, foreground='red')
                cnt = 0
                if self.title == "OO 계약" or not self.title:
                    message1.grid(row=2, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=2, column=4)
                if not self.pdct:
                    message2.grid(row=3, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=3, column=4)
                if not self.tests:
                    message3.grid(row=4, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=4, column=4)
                if not self.price.isnumeric():
                    message4.grid(row=5, column = 4)
                elif not int(self.price):
                    message4.grid(row=5, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=5, column=4)
                if not self.date_1.isnumeric() or len(self.date_1) != 8:
                    message5.grid(row=6, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=6, column=4)
                if not self.date_2.isnumeric() or len(self.date_2) != 8:
                    message6.grid(row=7, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=7, column=4)
                if not self.company:
                    message7.grid(row=8, column = 4)
                else:
                    cnt += 1
                    tk.Message(self.tkWindowContract, text="                    ", 
                                            width=100, padx = 10).grid(row=8, column=4)

                if cnt == 7:
                    if int(self.price) >= 1000000000000:
                        messagebox.showinfo("안내", '1조원이 넘는 금액은 입력할 수 없습니다.', parent=self.tkWindowContract)
                        pass
                    else:
                        date_contract = str(self.date_1)[:4]+"-"+str(self.date_1)[4:6]+"-"+str(self.date_1)[6:8]
                        date_end = str(self.date_2)[:4]+"-"+str(self.date_2)[4:6]+"-"+str(self.date_2)[6:8]
                        
                        #If all fields are correctly inputted, wrap all inputs in a dictionary for further use 
                        info = {
                            'title': self.title,
                            'pdct': self.pdct,
                            'tests': self.tests,
                            'price': self.price,
                            'date_contract': date_contract,
                            'date_end': date_end,
                            'company': self.company
                        }
                        if btn == 'con':
                            self.writeContract(info)
                        elif btn == 'chu':
                            self.writeChu(info)
                        else:
                            self.askName(info)

            conBtn1 = tk.Button(self.tkWindowContract, overrelief="solid", 
                               command = lambda: goContract(conEnTitle, conEnPdct, conEnCom,
                                                           conEnPrice, conEnStartDate,
                                                           conEnEndDate, selectedCompany, btn = 'con'), 
                               text = "표준계약서 작성")
            conBtn1.grid(row=9, column=2, pady = 35, sticky = 'e')
            
            conBtn2 = tk.Button(self.tkWindowContract, overrelief="solid", 
                               command = lambda: goContract(conEnTitle, conEnPdct, conEnCom,
                                                           conEnPrice, conEnStartDate,
                                                           conEnEndDate, selectedCompany, btn = 'chu'), 
                               text = "추진사유서 작성")
            conBtn2.grid(row=9, column=3)
            
            conBtn3 = tk.Button(self.tkWindowContract, overrelief="solid", width = 13,
                               command = lambda: goContract(conEnTitle, conEnPdct, conEnCom,
                                                           conEnPrice, conEnStartDate,
                                                           conEnEndDate, selectedCompany, btn = 'mail'), 
                               text = "품의 메일 양식")
            conBtn3.grid(row=9, column=4, sticky = 'w')
            
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}')    
    
    #Window for '프로그램 정보'
    def about(self):
        try:
            for child in self.window.winfo_children():
                if isinstance(child, tk.Toplevel):
                    if child.title() == '프로그램 정보':
                        child.destroy()

            tkWindowAbout = tk.Toplevel(self.window)
            tkWindowAbout.geometry(f"500x550+{screen_width//2 - 250}+{screen_height//2 - 100}")
            tkWindowAbout.title("프로그램 정보")
            tkWindowAbout.resizable(0, 0)
            tkWindowAbout.focus_force()

            self.img = tk.PhotoImage(file='./icon/icon_image.png')
            imgLb = tk.Label(tkWindowAbout, image=self.img)
            imgLb.pack(pady = (10, 0))

            aboutLb1 = tk.Label(tkWindowAbout, text = """\n업무 자동화 프로그램\n\n제작자: """)
            aboutLb1.pack(pady = (0, 10), padx = 30)

            aboutLb2 = tk.Label(tkWindowAbout, text = "Ver. 1.0.0 (2023.XX.XX)")
            aboutLb2.pack()
            
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.window)
 
    #Window for '설정'
    def config(self):
        try:
            for child in self.window.winfo_children():
                if isinstance(child, tk.Toplevel):
                    if child.title() == '설정':
                        print(child.title())
                        child.destroy()

            self.tkWindowConfig = tk.Toplevel(self.window)
            self.tkWindowConfig.geometry(f"800x600+{screen_width//2 + 100}+{screen_height//2 - 300}")
            self.tkWindowConfig.title("설정")
            self.tkWindowConfig.resizable(0, 0)
            self.tkWindowConfig.focus_force()

            notebookConfig = ttk.Notebook(self.tkWindowConfig)
            notebookConfig.pack(fill='both', expand=True)
            configFrame1 = tk.Frame(self.tkWindowConfig)
            notebookConfig.add(configFrame1, text="계약서 기본 정보")

            #Make an empty dictionary for entries
            crFr1Ens = {}
            tk.Label(configFrame1, text="").grid(row=1, pady = 10)
            
            cfFr1Lb_nameLG = tk.Label(configFrame1, text = "당사 상호")
            cfFr1Lb_nameLG.grid(row=2, column=2, pady = 10, padx = 10, sticky='w')
            cfFr1En_nameLG = tk.Entry(configFrame1)
            cfFr1En_nameLG.insert(0, data['common']['nameLG'])
            cfFr1En_nameLG.grid(row=2, column = 3)
            crFr1Ens[cfFr1En_nameLG] = (data['common']['nameLG'], 'nameLG')
            cfFr1Lb_nameLG2 = tk.Label(configFrame1, text = "당사 상호(서명부)")
            cfFr1Lb_nameLG2.grid(row=2, column=4, pady = 10, padx = 10, sticky='w')
            cfFr1En_nameLG2 = tk.Entry(configFrame1)
            cfFr1En_nameLG2.insert(0, data['common']['nameLG2'])
            cfFr1En_nameLG2.grid(row=2, column = 5, sticky = 'w')
            crFr1Ens[cfFr1En_nameLG2] = (data['common']['nameLG2'], 'nameLG2')
            cfFr1Lb_dayPayment = tk.Label(configFrame1, text = "선금/잔금 지급일")
            cfFr1Lb_dayPayment.grid(row=2, column=6, pady = 10, padx = 10, sticky='w')
            cfFr1En_dayPayment = tk.Entry(configFrame1, width = 12)
            cfFr1En_dayPayment.insert(0, data['common']['dayPayment'])
            cfFr1En_dayPayment.grid(row=2, column = 7, sticky = 'w')
            crFr1Ens[cfFr1En_dayPayment] = (data['common']['dayPayment'], 'dayPayment')
            
            cfFr1Lb_ceoLG = tk.Label(configFrame1, text = "당사 대표이사명")
            cfFr1Lb_ceoLG.grid(row=3, column=2, pady = 10, padx = 10, sticky='w')
            cfFr1En_ceoLG = tk.Entry(configFrame1)
            cfFr1En_ceoLG.insert(0, data['common']['ceoLG'])
            cfFr1En_ceoLG.grid(row=3, column = 3)
            crFr1Ens[cfFr1En_ceoLG] = (data['common']['ceoLG'], 'ceoLG')
            cfFr1Lb_businessLG = tk.Label(configFrame1, text = "당사 업태 및 종목")
            cfFr1Lb_businessLG.grid(row=3, column=4, pady = 10, padx = 10, sticky='w')
            cfFr1En_businessLG = tk.Entry(configFrame1, width = 53)
            cfFr1En_businessLG.insert(0, data['common']['businessLG'])
            cfFr1En_businessLG.grid(row=3, column = 5, columnspan=3, sticky = 'w')
            crFr1Ens[cfFr1En_businessLG] = (data['common']['businessLG'], 'businessLG')
            
            cfFr1Lb_regNumLG = tk.Label(configFrame1, text = "당사 사업자번호")
            cfFr1Lb_regNumLG.grid(row=4, column=2, pady = 10, padx = 10, sticky='w')
            cfFr1En_regNumLG = tk.Entry(configFrame1)
            cfFr1En_regNumLG.insert(0, data['common']['regNumLG'])
            cfFr1En_regNumLG.grid(row=4, column = 3)
            crFr1Ens[cfFr1En_regNumLG] = (data['common']['regNumLG'], 'regNumLG')
            cfFr1Lb_adressLG = tk.Label(configFrame1, text = "당사 주소")
            cfFr1Lb_adressLG.grid(row=4, column=4, pady = 10, padx = 10, sticky='w')
            cfFr1En_adressLG = tk.Entry(configFrame1, width = 53)
            cfFr1En_adressLG.insert(0, data['common']['adressLG'])
            cfFr1En_adressLG.grid(row=4, column = 5, columnspan=3, sticky = 'w')
            crFr1Ens[cfFr1En_adressLG] = (data['common']['adressLG'], 'adressLG')
            
            tk.Label(configFrame1, text = "").grid(row=5, column=1)
            
            cfFr1Lb_personLG = tk.Label(configFrame1, text = "의뢰 책임자")
            cfFr1Lb_personLG.grid(row=6, column=2, pady = 10, padx = 10, sticky='w')
            cfFr1En_personLG = tk.Entry(configFrame1)
            cfFr1En_personLG.insert(0, data['common']['personLG'])
            cfFr1En_personLG.grid(row=6, column = 3)
            crFr1Ens[cfFr1En_personLG] = (data['common']['personLG'], 'personLG')
            cfFr1Lb_personTeamLG = tk.Label(configFrame1, text = "책임자 소속")
            cfFr1Lb_personTeamLG.grid(row=6, column=4, pady = 10, padx = 10, sticky='w')
            cfFr1En_personTeamLG = tk.Entry(configFrame1, width = 23)
            cfFr1En_personTeamLG.insert(0, data['common']['personTeamLG'])
            cfFr1En_personTeamLG.grid(row=6, column = 5)
            crFr1Ens[cfFr1En_personTeamLG] = (data['common']['personTeamLG'], 'personTeamLG')
            cfFr1Lb_personContactLG = tk.Label(configFrame1, text = "책임자 연락처")
            cfFr1Lb_personContactLG.grid(row=6, column=6, pady = 10, padx = 10)
            cfFr1En_personContactLG = tk.Entry(configFrame1, width = 12)
            cfFr1En_personContactLG.insert(0, data['common']['personContactLG'])
            cfFr1En_personContactLG.grid(row=6, column = 7, sticky = 'w')
            crFr1Ens[cfFr1En_personContactLG] = (data['common']['personContactLG'], 'personContactLG')
        
            cfFr1Lb_personEmailLG = tk.Label(configFrame1, text = "책임자 이메일")
            cfFr1Lb_personEmailLG.grid(row=7, column=2, pady = 10, padx = 10, sticky='w')
            cfFr1En_personEmailLG = tk.Entry(configFrame1)
            cfFr1En_personEmailLG.insert(0, data['common']['personEmailLG'])
            cfFr1En_personEmailLG.grid(row=7, column = 3)
            crFr1Ens[cfFr1En_personEmailLG] = (data['common']['personEmailLG'], 'personEmailLG')
            cfFr1Lb_personAddressLG = tk.Label(configFrame1, text = "책임자 주소")
            cfFr1Lb_personAddressLG.grid(row=7, column=4, pady = 10, padx = 10, sticky='w')
            cfFr1En_personAddressLG = tk.Entry(configFrame1, width = 53)
            cfFr1En_personAddressLG.insert(0, data['common']['personAddressLG'])
            cfFr1En_personAddressLG.grid(row=7, column = 5, columnspan=3, sticky = 'w')
            crFr1Ens[cfFr1En_personAddressLG] = (data['common']['personAddressLG'], 'personAddressLG')
            
            tk.Label(configFrame1, text = "").grid(row=8, column=1)
            
            cfFr1Lb_howToUse = tk.Label(configFrame1, text = "제품 정보")
            cfFr1Lb_howToUse.grid(row=9, column=2, pady = 10, padx = 10, sticky='w')
            cfFr1En_howToUse = tk.Entry(configFrame1)
            cfFr1En_howToUse.insert(0, data['common']['howToUse'])
            cfFr1En_howToUse.grid(row=9, column = 3)
            crFr1Ens[cfFr1En_howToUse] = (data['common']['howToUse'], 'howToUse')
            cfFr1Lb_howToStore = tk.Label(configFrame1, text = "보관 방법")
            cfFr1Lb_howToStore.grid(row=9, column=4, pady = 10, padx = 10, sticky='w')
            cfFr1En_howToStore = tk.Entry(configFrame1, width = 53)
            cfFr1En_howToStore.insert(0, data['common']['howToStore'])
            cfFr1En_howToStore.grid(row=9, column = 5, columnspan=3, sticky = 'w')
            crFr1Ens[cfFr1En_howToStore] = (data['common']['howToStore'], 'howToStore')

            #Make entries readonly
            for entry in crFr1Ens:
                entry.config(state = 'disabled')
                        
            def startConfig1():
                cfFr1Btn1.config(state = "disabled")
                cfFr1Btn3.config(state = "disabled")
                for entry in crFr1Ens:
                    entry.config(state = 'normal')
                cfFr1Btn2.grid(row=10, column=5, padx = 20, sticky = 'w')
       
            def finishConfig1():
                #Check if there is any change first
                if [entry.get() for entry in crFr1Ens] != [i[0] for i in list(crFr1Ens.values())]:

                    askConfig1 = messagebox.askyesnocancel('계약서 기본 정보 변경', '표준계약서 기본 정보를 변경하시겠어요?\n(작업 중인 창은 종료됩니다)',
                                                    parent=self.tkWindowConfig)
                    if askConfig1 is None:
                        self.config()
                    elif askConfig1:
                        for entry in crFr1Ens:
                            if entry.get() != crFr1Ens[entry][0]:
                                tempKey = crFr1Ens[entry][1]
                                data['common'][tempKey] = entry.get()
                                crFr1Ens[entry] = (data['common'][tempKey], tempKey)
                            entry.config(state = 'disabled')

                        saveData(data)
                        loadData()
                        
                        cfFr1Btn1.config(state = "normal")
                        cfFr1Btn3.config(state = "normal")
                        cfFr1Btn2.grid_remove()
                        for child in self.window.winfo_children():
                            if isinstance(child, tk.Toplevel):
                                child.destroy()
                        messagebox.showinfo('계약서 기본 정보 변경', '표준계약서 기본 정보가 변경되었습니다')
                    else:
                        pass
                else:
                    messagebox.showwarning('계약서 기본 정보 변경', '수정하신 내용이 없어요. 내용 수정 후 확인을 눌러주세요.', parent=self.tkWindowConfig)
                        
            #Reset configuration using data_default.json    
            def resetConfig1():
                askResetConfig1 = messagebox.askyesno('계약서 기본 정보 초기화', 
                                                 '''표준계약서 기본 정보가 초기화됩니다!! \n진행하시겠어요? (작업 중인 창은 종료됩니다)''',
                                                      icon = 'warning',
                                parent=self.tkWindowConfig)
                if askResetConfig1:
                    with open('./pmdata/data_default.json', 'r') as file:
                        data_default = json.load(file)
                        
                    saveData(data_default)
                    loadData()
                    
                    for child in self.window.winfo_children():
                        if isinstance(child, tk.Toplevel):
                            child.destroy()
                    messagebox.showinfo('계약서 기본 정보 초기화', '표준계약서 기본 정보가 초기화되었습니다.')
                    
                else:
                    pass
                
            cfFr1Btn1 = tk.Button(configFrame1, overrelief="solid", text = "수정", command = startConfig1,
                                 width = 10)
            cfFr1Btn1.grid(row=10, column=4, pady = 20)
            cfFr1Btn2 = tk.Button(configFrame1, overrelief="solid", text = "확인", width = 10, 
                                  command = finishConfig1, foreground="red")
            cfFr1Btn3 = tk.Button(configFrame1, overrelief="solid", text = "초기화", width = 10, 
                          command = resetConfig1)
            cfFr1Btn3.grid(row=10, column=7, pady = 20, sticky = 'e')

        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}', parent=self.tkWindowConfig)
        

    def __main__(self):
        try:
            menubar=tk.Menu(self.window)

            menu_1=tk.Menu(menubar, tearoff=0)
            menu_1.add_command(label="설정", command=self.config)
            menu_1.add_command(label="프로그램 정보", command=self.about)
            menu_1.add_separator()
            menu_1.add_command(label="종료", command=self.close)
            menubar.add_cascade(label="메뉴", menu=menu_1)
            self.window.config(menu=menubar)
            
            mainLb = tk.Label(self.window, text = "")
            mainLb.pack(side = 'top',  pady = 70, padx = 30)
            tkBtnContract = tk.Button(self.window, text="기본 연구 업무",
                                          command = self.startContract)
            tkBtnContract.pack(anchor = 'center', ipadx = 10, ipady = 4)
        except Exception as e:
            messagebox.showinfo("Error", f'에러가 생겼어요. 다시 한 번 더 해보시고,\n
                반복되면 아래 에러 내용을 제작자에게 보내주세요.\n\n{str(e)}')

if __name__ == '__main__':    
    window = tk.Tk()
    screen_width = window.winfo_screenwidth()
    screen_height = window.winfo_screenheight()
    windows_tkinter(window)
    window.mainloop()
